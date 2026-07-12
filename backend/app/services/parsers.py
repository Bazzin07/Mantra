import csv
import io
import os
from typing import List

from pydantic import BaseModel


class UnsupportedFormatError(ValueError):
    pass


class DependencyUnavailableError(RuntimeError):
    pass


class DocumentParseError(ValueError):
    pass


class ParsedDocument(BaseModel):
    document_type: str
    pages: List[str]


class DocumentParser:
    def parse(self, filename: str, content_type: str, payload: bytes) -> ParsedDocument:
        extension = os.path.splitext(filename.lower())[1]
        if extension in {".txt", ".md", ".eml"}:
            return ParsedDocument(document_type=extension.lstrip("."), pages=[decode_text(payload)])
        if extension == ".csv":
            return ParsedDocument(document_type="csv", pages=[parse_csv(payload)])
        if extension == ".pdf":
            return parse_pdf(payload)
        if extension == ".docx":
            return parse_docx(payload)
        if extension == ".pptx":
            return parse_pptx(payload)
        if extension in {".xlsx", ".xls"}:
            return parse_spreadsheet(payload)
        if extension in {".png", ".jpg", ".jpeg", ".tif", ".tiff"}:
            return parse_image(payload)
        raise UnsupportedFormatError(f"Unsupported document format: {extension or content_type}")


def decode_text(payload: bytes) -> str:
    for encoding in ("utf-8", "utf-16", "latin-1"):
        try:
            return payload.decode(encoding)
        except UnicodeDecodeError:
            continue
    raise ValueError("Could not decode text payload")


def parse_csv(payload: bytes) -> str:
    text = decode_text(payload)
    reader = csv.reader(io.StringIO(text))
    rows = [" | ".join(cell.strip() for cell in row) for row in reader]
    return "\n".join(rows)


def parse_pdf(payload: bytes) -> ParsedDocument:
    try:
        import fitz
    except ImportError as exc:
        raise DependencyUnavailableError("PDF parsing requires PyMuPDF") from exc

    pages: List[str] = []
    try:
        with fitz.open(stream=payload, filetype="pdf") as document:
            for page in document:
                pages.append(page.get_text("text"))
    except Exception as exc:
        raise DocumentParseError("Could not parse PDF document") from exc
    return ParsedDocument(document_type="pdf", pages=pages)


def parse_docx(payload: bytes) -> ParsedDocument:
    try:
        from docx import Document
    except ImportError as exc:
        raise DependencyUnavailableError("DOCX parsing requires python-docx") from exc

    try:
        document = Document(io.BytesIO(payload))
    except Exception as exc:
        raise DocumentParseError("Could not parse DOCX document") from exc
    paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            paragraphs.append(" | ".join(cell.text.strip() for cell in row.cells))
    return ParsedDocument(document_type="docx", pages=["\n".join(paragraphs)])


def parse_pptx(payload: bytes) -> ParsedDocument:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise DependencyUnavailableError("PPTX parsing requires python-pptx") from exc

    try:
        presentation = Presentation(io.BytesIO(payload))
    except Exception as exc:
        raise DocumentParseError("Could not parse PPTX document") from exc

    pages: List[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        fragments = [f"Slide {index}"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                fragments.append(shape.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    fragments.append(" | ".join(cell.text.strip() for cell in row.cells))
        pages.append("\n".join(fragments))
    return ParsedDocument(document_type="pptx", pages=pages)


def parse_spreadsheet(payload: bytes) -> ParsedDocument:
    try:
        import openpyxl
    except ImportError as exc:
        raise DependencyUnavailableError("Spreadsheet parsing requires openpyxl") from exc

    try:
        workbook = openpyxl.load_workbook(io.BytesIO(payload), data_only=True)
    except Exception as exc:
        raise DocumentParseError("Could not parse spreadsheet document") from exc
    pages: List[str] = []
    for sheet in workbook.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            rows.append(" | ".join("" if value is None else str(value) for value in row))
        pages.append(f"Sheet: {sheet.title}\n" + "\n".join(rows))
    return ParsedDocument(document_type="spreadsheet", pages=pages)


def parse_image(payload: bytes) -> ParsedDocument:
    try:
        from PIL import Image
        import pytesseract
    except ImportError as exc:
        raise DependencyUnavailableError("Image OCR requires Pillow and pytesseract") from exc

    try:
        image = Image.open(io.BytesIO(payload))
        text = pytesseract.image_to_string(image)
    except Exception as exc:
        raise DocumentParseError("Could not OCR image document") from exc
    return ParsedDocument(document_type="image", pages=[text])
